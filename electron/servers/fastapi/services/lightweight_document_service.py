"""
Lightweight document converter for Windows/MSIX compatibility.
Uses pure-Python libraries: pdfplumber for PDF, docx2txt for DOCX, python-pptx for PPTX.
No subprocess, no external runtimes, MSIX/Appx safe.
"""
import os
from typing import List, Optional

import docx2everything
import pdfplumber
from pptx import Presentation


class LightweightDocumentConverter:
    """Lightweight document converter supporting PDF, DOCX, and PPTX."""
    
    def convert(self, file_path: str) -> str:
        """
        Convert document to markdown text.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text in markdown format
            
        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file does not exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self._convert_pdf(file_path)
        elif file_ext == '.docx':
            return self._convert_docx(file_path)
        elif file_ext == '.pptx':
            return self._convert_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _convert_pdf(self, path: str) -> str:
        """
        Convert PDF to markdown using pdfplumber.

        Args:
            path: Path to PDF file

        Returns:
            Extracted text in markdown format
        """
        texts: List[str] = []
        with pdfplumber.open(path) as pdf:
            for idx, page in enumerate(pdf.pages):
                page_text = f"## Page {idx + 1}\n"
                page_text += page.extract_text() or ""
                texts.append(page_text)
        return "\n\n".join(texts)
    
    def _convert_docx(self, path: str) -> str:
        """
        Extract markdown from DOCX using docx2everything (no images).

        Args:
            path: Path to DOCX file

        Returns:
            Extracted markdown (no images)
        """
        # Use the correct API: process_to_markdown(path) without img_dir extracts markdown without images
        markdown = docx2everything.process_to_markdown(path)
        return markdown if markdown else ""
    
    def _convert_pptx(self, path: str) -> str:
        """
        Convert PPTX to markdown using python-pptx.
        
        Args:
            path: Path to PPTX file
            
        Returns:
            Extracted text in markdown format
        """
        prs = Presentation(path)
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            
            # Extract slide title (usually first shape with title placeholder)
            title_text = None
            for shape in slide.shapes:
                if hasattr(shape, "placeholder"):
                    if shape.placeholder.placeholder_format.type == 1:  # Title placeholder
                        if hasattr(shape, "text") and shape.text.strip():
                            title_text = shape.text.strip()
                            break
            
            # If no title placeholder found, try to find text box at top
            if not title_text:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's likely a title (first text shape, short text)
                        text = shape.text.strip()
                        if len(text) < 200:  # Heuristic: titles are usually short
                            title_text = text
                            break
            
            # Add slide title
            if title_text:
                slide_parts.append(f"# {title_text}")
            else:
                slide_parts.append(f"# Slide {slide_num}")
            
            # Extract content (bullet points and text)
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                
                text = shape.text.strip()
                if not text:
                    continue
                
                # Skip if this is the title we already added
                if title_text and text == title_text:
                    continue
                
                # Check if it's a text frame with paragraphs (bullet points)
                if hasattr(shape, "text_frame"):
                    paragraphs = shape.text_frame.paragraphs
                    if len(paragraphs) > 1:
                        # Multiple paragraphs - likely bullet points
                        for para in paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                # Check bullet level
                                level = para.level
                                indent = "  " * level
                                slide_parts.append(f"{indent}- {para_text}")
                    else:
                        # Single paragraph
                        if text and text != title_text:
                            slide_parts.append(text)
                else:
                    # Plain text shape
                    if text and text != title_text:
                        slide_parts.append(text)
            
            if slide_parts:
                markdown_parts.append("\n".join(slide_parts))
        
        return "\n\n---\n\n".join(markdown_parts)


class DocumentService:
    """
    Document service wrapper providing parse_to_markdown interface.
    Same parse_to_markdown entry point as LiteParseService for optional Windows fallback.
    """
    
    def __init__(self):
        self.converter = LightweightDocumentConverter()
    
    def parse_to_markdown(self, file_path: str) -> str:
        """
        Parse document to markdown format.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text in markdown format
        """
        return self.converter.convert(file_path)
